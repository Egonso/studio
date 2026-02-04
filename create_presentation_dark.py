import re
import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Configuration
MD_FILE = "/Users/momofeichtinger/.gemini/antigravity/brain/4e860bf8-e092-4a74-815c-e716d3a67362/presentation_outline.md"
OUTPUT_FILE = "/Users/momofeichtinger/.gemini/antigravity/brain/4e860bf8-e092-4a74-815c-e716d3a67362/Eukigesetz_Studio_Presentation_Dark.pptx"
ASSETS_DIR = "/Users/momofeichtinger/.gemini/antigravity/brain/4e860bf8-e092-4a74-815c-e716d3a67362"

# Dark Mode Colors
DARK_BG = RGBColor(10, 20, 40)   # Very Dark Navy/Black
HEADER_BG = RGBColor(0, 45, 179) # Deep Royal Blue (Branding)
WHITE = RGBColor(255, 255, 255)  # Define White explicitly
TEXT_TITLE = WHITE 
TEXT_BODY = RGBColor(220, 220, 230)  # Light Grey/White-ish
ACCENT_LINE = RGBColor(0, 150, 255)  # Cyan/Bright Blue for separation lines

# Image Mapping
IMAGE_MAP = {
    1: "slide_01_title_background_*.png",
    3: "slide_03_problem_complexity_*.png",
    4: "slide_04_dashboard_mockup_*.png",
    7: "slide_07_structured_process_*.png",
    13: "slide_13_trust_portal_*.png",
    20: "slide_20_value_stack_*.png",
    25: "slide_25_call_to_action_*.png"
}

def find_image(slide_num):
    pattern = IMAGE_MAP.get(slide_num)
    if not pattern:
        return None
    files = glob.glob(os.path.join(ASSETS_DIR, pattern))
    if files:
        files.sort(key=os.path.getmtime, reverse=True)
        return files[0]
    return None

def clean_md_field(line):
    if ':' in line:
        return line.split(':', 1)[1].strip()
    return line

def parse_markdown(md_path):
    with open(md_path, 'r') as f:
        content = f.read()
    
    slides = []
    current_slide = None
    
    lines = content.split('\n')
    for line in lines:
        line = line.strip()
        if not line: continue
        
        if line.startswith('## Slide'):
            if current_slide:
                slides.append(current_slide)
            current_slide = {'title': '', 'subtitle': '', 'content': [], 'visual_note': ''}
        elif current_slide is not None:
            if line.startswith('**Titel:**') or line.startswith('**Title:**'):
                 current_slide['title'] = clean_md_field(line)
            elif line.startswith('**Untertitel:**') or line.startswith('**Subtitle:**'):
                 current_slide['subtitle'] = clean_md_field(line)
            elif line.startswith('**Text:**'):
                 current_slide['content'].append(clean_md_field(line))
            elif line.startswith('**Visual:**'):
                 current_slide['visual_note'] = clean_md_field(line)
            elif not current_slide['title'] and line.startswith('#'):
                 current_slide['title'] = line.lstrip('#').strip()

    if current_slide:
        slides.append(current_slide)
        
    return slides

def render_markdown_text(paragraph, text, font_size=Pt(18), color=TEXT_BODY):
    """
    Parses text for **bold** and strictly removes all * characters from output.
    """
    parts = re.split(r'(\*\*.*?\*\*)', text)
    
    for part in parts:
        if not part: continue
        
        run = paragraph.add_run()
        run.font.size = font_size
        run.font.color.rgb = color
        
        if part.startswith('**') and part.endswith('**'):
            clean_text = part[2:-2] # Remove ** marker
            clean_text = clean_text.replace('*', '') 
            run.text = clean_text
            run.font.bold = True
        else:
            clean_text = part.replace('*', '')
            run.text = clean_text
            run.font.bold = False

def create_pptx(slides):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides):
        slide_num = i + 1
        img_path = find_image(slide_num)
        
        if slide_num == 1:
            # === Title Slide ===
            slide = prs.slides.add_slide(prs.slide_layouts[6]) 
            
            # Dark Background Image or Color
            if img_path:
                slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            else:
                 bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
                 bg.fill.solid()
                 bg.fill.fore_color.rgb = DARK_BG
            
            # Text Setup
            top = Inches(2.0)
            left = Inches(1)
            width = Inches(11.333)
            height = Inches(2.5)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            description = slide_data['title']
            p = txBox.text_frame.add_paragraph()
            p.text = description.replace('*', '') 
            p.font.bold = True
            p.font.size = Pt(44)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.LEFT
            
            if slide_data.get('subtitle'):
                top_sub = top + Inches(1.5)
                txBoxSub = slide.shapes.add_textbox(left, top_sub, width, height)
                p_sub = txBoxSub.text_frame.add_paragraph()
                p_sub.text = slide_data['subtitle'].replace('*', '')
                p_sub.font.size = Pt(24)
                p_sub.font.color.rgb = WHITE
                p_sub.alignment = PP_ALIGN.LEFT

        else:
            # === Content Slide (Dark Mode) ===
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # 1. Background Fill (Dark)
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = DARK_BG
            bg.line.fill.background()
            
            # 2. Header Bar (Branding Blue)
            header_height = Inches(1.2)
            header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, header_height)
            header.fill.solid()
            header.fill.fore_color.rgb = HEADER_BG
            header.line.fill.background()
            
            # 3. Accent Line (Cyan separator)
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, header_height, prs.slide_width, Inches(0.05))
            line.fill.solid()
            line.fill.fore_color.rgb = ACCENT_LINE
            line.line.fill.background()

            # 4. Slide Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0), Inches(12.333), header_height)
            tf_title = title_box.text_frame
            tf_title.vertical_anchor = MSO_ANCHOR.MIDDLE
            p_title = tf_title.add_paragraph()
            p_title.text = slide_data['title'].replace('*', '')
            p_title.font.bold = True
            p_title.font.size = Pt(32)
            p_title.font.color.rgb = TEXT_TITLE 
            p_title.alignment = PP_ALIGN.LEFT
            
            content_top = Inches(1.7) 
            
            if img_path:
                # Two Column Layout
                pic_left = Inches(7.5)
                pic_top = content_top
                pic_width = Inches(5.3)
                
                # Image Border for design pop
                border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, pic_left - Inches(0.05), pic_top - Inches(0.05), pic_width + Inches(0.1), (pic_width*0.75) + Inches(0.1)) # Approx height
                border.fill.solid()
                border.fill.fore_color.rgb = ACCENT_LINE
                border.line.fill.background()
                
                slide.shapes.add_picture(img_path, pic_left, pic_top, width=pic_width)
                
                text_width = Inches(6.5)
            else:
                text_width = Inches(12)
            
            # Content Box
            if slide_data['content']:
                txBox = slide.shapes.add_textbox(Inches(0.5), content_top, text_width, Inches(5.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                
                for line in slide_data['content']:
                    p = tf.add_paragraph()
                    render_markdown_text(p, "• " + line, font_size=Pt(20), color=TEXT_BODY)
                    p.space_after = Pt(14)

    prs.save(OUTPUT_FILE)
    print(f"Presentation saved to {OUTPUT_FILE}")

if __name__ == "__main__":
    slides = parse_markdown(MD_FILE)
    print(f"Found {len(slides)} slides")
    create_pptx(slides)
