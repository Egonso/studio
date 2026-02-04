import re
import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Configuration
MD_FILE = "/Users/momofeichtinger/.gemini/antigravity/brain/4e860bf8-e092-4a74-815c-e716d3a67362/presentation_outline.md"
OUTPUT_FILE = "/Users/momofeichtinger/.gemini/antigravity/brain/4e860bf8-e092-4a74-815c-e716d3a67362/Eukigesetz_Studio_Presentation.pptx"
ASSETS_DIR = "/Users/momofeichtinger/.gemini/antigravity/brain/4e860bf8-e092-4a74-815c-e716d3a67362"

# Colors
BLUE = RGBColor(0, 45, 179)      # Deep Royal Blue
LIGHT_BLUE = RGBColor(230, 240, 255) # Light nuance
WHITE = RGBColor(255, 255, 255)
GREY = RGBColor(80, 80, 80)
BLACK = RGBColor(0, 0, 0)

# Image Mapping
IMAGE_MAP = {
    1: "slide_01_title_background_*.png",
    3: "slide_03_problem_complexity_*.png",
    4: "slide_04_dashboard_mockup_*.png",
    7: "slide_07_structured_process_*.png",
    13: "slide_13_trust_portal_*.png",
    20: "slide_20_value_stack_*.png",
    25: "slide_25_call_to_action_*.png"
}

def find_image(slide_num):
    pattern = IMAGE_MAP.get(slide_num)
    if not pattern:
        return None
    files = glob.glob(os.path.join(ASSETS_DIR, pattern))
    if files:
        # Sort by modification time to get the latest
        files.sort(key=os.path.getmtime, reverse=True)
        return files[0]
    return None

def clean_md_field(line):
    # Removes the "**Key:** " part
    if ':' in line:
        return line.split(':', 1)[1].strip()
    return line

def parse_markdown(md_path):
    with open(md_path, 'r') as f:
        content = f.read()
    
    slides = []
    current_slide = None
    
    lines = content.split('\n')
    for line in lines:
        line = line.strip()
        if not line: continue
        
        if line.startswith('## Slide'):
            if current_slide:
                slides.append(current_slide)
            current_slide = {'title': '', 'subtitle': '', 'content': [], 'visual_note': ''}
        elif current_slide is not None:
            if line.startswith('**Titel:**') or line.startswith('**Title:**'):
                 current_slide['title'] = clean_md_field(line)
            elif line.startswith('**Untertitel:**') or line.startswith('**Subtitle:**'):
                 current_slide['subtitle'] = clean_md_field(line)
            elif line.startswith('**Text:**'):
                 current_slide['content'].append(clean_md_field(line))
            elif line.startswith('**Visual:**'):
                 current_slide['visual_note'] = clean_md_field(line)
            # Handle implicit title if missing and it's a bold line
            elif not current_slide['title'] and line.startswith('#'):
                 current_slide['title'] = line.lstrip('#').strip()

    if current_slide:
        slides.append(current_slide)
        
    return slides

def render_markdown_text(paragraph, text, font_size=Pt(18), color=GREY):
    """
    Parses text for **bold** markers and creates runs.
    Ensures all ** are removed from output.
    """
    parts = re.split(r'(\*\*.*?\*\*)', text)
    
    for part in parts:
        if not part: continue
        
        run = paragraph.add_run()
        run.font.size = font_size
        run.font.color.rgb = color
        
        if part.startswith('**') and part.endswith('**'):
            clean_text = part[2:-2] # Remove **
            run.text = clean_text
            run.font.bold = True
        else:
            # Safety cleanup: if loose ** exist, remove them
            run.text = part.replace('**', '') 
            run.font.bold = False

def create_pptx(slides):
    prs = Presentation()
    
    # Set slide dimensions (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides):
        slide_num = i + 1
        img_path = find_image(slide_num)
        
        if slide_num == 1:
            # === Title Slide (Reverted to V1 Style) ===
            slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
            
            # Background Image (Full Slide)
            if img_path:
                slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            else:
                 # Fallback Blue background
                 bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
                 bg.fill.solid()
                 bg.fill.fore_color.rgb = BLUE
            
            # Title Box (On top of image, Left Aligned)
            # Move text down a bit to match V1 feel (usually centered or slightly lower)
            top = Inches(2.0)
            left = Inches(1)
            width = Inches(11.333)
            height = Inches(2.5)
            
            # Add Title
            txBox = slide.shapes.add_textbox(left, top, width, height)
            description = slide_data['title']
            p = txBox.text_frame.add_paragraph()
            p.text = description
            p.font.bold = True
            p.font.size = Pt(44)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.LEFT
            
            # Add Subtitle
            if slide_data.get('subtitle'):
                top_sub = top + Inches(1.5) # Spacing below title
                txBoxSub = slide.shapes.add_textbox(left, top_sub, width, height)
                p_sub = txBoxSub.text_frame.add_paragraph()
                p_sub.text = slide_data['subtitle']
                p_sub.font.size = Pt(24)
                p_sub.font.color.rgb = WHITE
                p_sub.alignment = PP_ALIGN.LEFT

        else:
            # === Content Slide ===
            slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
            
            # Header Bar (Branding) - Blue
            header_height = Inches(1.2)
            header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, header_height)
            header.fill.solid()
            header.fill.fore_color.rgb = BLUE
            header.line.fill.background()
            
            # Slide Title - IN THE BLUE BAR
            # Centered vertically in the header
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0), Inches(12.333), header_height)
            tf_title = title_box.text_frame
            tf_title.vertical_anchor = MSO_ANCHOR.MIDDLE # Center vertically in the blue bar
            p_title = tf_title.add_paragraph()
            p_title.text = slide_data['title']
            p_title.font.bold = True
            p_title.font.size = Pt(32)
            p_title.font.color.rgb = WHITE
            p_title.alignment = PP_ALIGN.LEFT
            
            content_top = Inches(1.6) # Spacing below header
            
            if img_path:
                # Two Column Layout
                # Image Right
                pic_left = Inches(7.5)
                pic_top = content_top
                pic_width = Inches(5.3)
                
                # Add picture
                slide.shapes.add_picture(img_path, pic_left, pic_top, width=pic_width)
                
                # Text Left
                text_width = Inches(6.5)
            else:
                # Full Width Text
                text_width = Inches(12)
            
            # Content Box
            if slide_data['content']:
                txBox = slide.shapes.add_textbox(Inches(0.5), content_top, text_width, Inches(5.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                
                for line in slide_data['content']:
                    p = tf.add_paragraph()
                    render_markdown_text(p, "• " + line, font_size=Pt(20), color=GREY)
                    p.space_after = Pt(12)

    prs.save(OUTPUT_FILE)
    print(f"Presentation saved to {OUTPUT_FILE}")

if __name__ == "__main__":
    slides = parse_markdown(MD_FILE)
    print(f"Found {len(slides)} slides")
    create_pptx(slides)
